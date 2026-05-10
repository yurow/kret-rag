"""
文档处理服务
负责文档上传、解析、分块、向量化
"""
from typing import List, Dict, Any, Optional
from fastapi import UploadFile
import uuid
import os
import re
import logging
from datetime import datetime
from pathlib import Path

from app.models.schemas import (
    DocumentUploadRequest, 
    DocumentResponse, 
    DocumentChunk,
    DocumentStatus,
    DocumentUploadResult
)
from app.core.config import settings
from app.db.database import db_manager
from app.db.db_migrator import run_migrations
from app.repositories.document_repository import DocumentRepository
from app.services.vector_service import vector_store_service

# 配置日志
logger = logging.getLogger(__name__)


class DocumentService:
    """文档处理服务"""
    
    def __init__(self):
        """初始化文档服务，创建上传目录和数据库表"""
        self.upload_dir = Path(settings.UPLOAD_DIR)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        # 创建清洗后文本存储目录
        self.text_upload_dir = Path("./uploads_text")
        self.text_upload_dir.mkdir(parents=True, exist_ok=True)
        
        # 确保数据目录存在
        data_dir = Path("./data")
        data_dir.mkdir(parents=True, exist_ok=True)
        
        # 执行数据库迁移（在创建表之前）
        logger.info("检查并执行数据库迁移...")
        run_migrations()
        
        # 创建数据库表
        db_manager.create_tables()
        logger.info("数据库初始化完成")
        
        # ⭐ 向量服务将在首次使用时异步初始化（避免阻塞启动）
        self._vector_service_initialized = False
    
    async def _ensure_vector_service(self):
        """确保向量服务已初始化（懒加载）"""
        if not self._vector_service_initialized:
            try:
                await vector_store_service.initialize()
                self._vector_service_initialized = True
                logger.info("向量服务初始化成功")
            except Exception as e:
                logger.error(f"向量服务初始化失败: {str(e)}", exc_info=True)
                raise

    async def upload_document(
        self, 
        file: UploadFile, 
        metadata: Optional[Dict[str, Any]] = None
    ) -> DocumentUploadResult:
        """
        上传并处理文档
        
        Args:
            file: 上传的文件对象
            metadata: 可选的元数据字典
            
        Returns:
            DocumentUploadResult: 包含文档响应和重复标识的结果对象
            
        Raises:
            ValueError: 当文件格式不支持、大小超限或文件已存在时
        """
        logger.info(f"开始处理文件上传: {file.filename}, 大小: {file.size} bytes")
        
        # 验证文件大小
        if file.size and file.size > settings.MAX_DOCUMENT_SIZE:
            logger.warning(f"文件大小超出限制: {file.size} > {settings.MAX_DOCUMENT_SIZE}")
            raise ValueError(f"File size exceeds maximum limit of {settings.MAX_DOCUMENT_SIZE}")
        
        # 验证文件格式
        file_extension = file.filename.split('.')[-1].lower()
        if file_extension not in settings.SUPPORTED_FORMATS:
            logger.warning(f"不支持的文件格式: {file_extension}")
            raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: {', '.join(settings.SUPPORTED_FORMATS)}")
        
        logger.info(f"文件格式验证通过: {file_extension}")
        
        db_session = db_manager.SessionLocal()
        
        try:
            # 创建仓储实例
            repo = DocumentRepository(db_session)
            
            # 检查是否存在重复文件（同名 + 同大小）
            logger.info(f"检查是否存在重复文件: {file.filename}, 大小: {file.size}")
            existing_doc = repo.find_duplicate(file.filename, file.size)
            if existing_doc and False:
                logger.info(f"发现重复文件，直接返回已有文档: {existing_doc.document_id}")
                # 文件已存在，直接返回已有文档信息
                doc_response = DocumentResponse(
                    document_id=existing_doc.document_id,
                    file_name=existing_doc.file_name,
                    status=DocumentStatus(existing_doc.status),
                    created_at=existing_doc.created_at,
                    updated_at=existing_doc.updated_at,
                    metadata=existing_doc.extra_metadata
                )
                return DocumentUploadResult(
                    document_response=doc_response,
                    is_duplicate=True
                )
            
            # 生成文档ID和文件名
            document_id = str(uuid.uuid4())
            safe_filename = f"{document_id}_{file.filename}"
            file_path = self.upload_dir / safe_filename
            
            logger.info(f"生成文档ID: {document_id}, 保存路径: {file_path}")
            
            # 保存文件到服务器
            logger.info("开始保存文件到上传目录...")
            contents = await file.read()
            with open(file_path, "wb") as f:
                f.write(contents)
            logger.info(f"文件保存成功: {file_path}, 大小: {len(contents)} bytes")
            
            # 提取文本内容
            logger.info(f"开始提取文本内容，文件格式: {file_extension}")
            extracted_text = await self.extract_text(str(file_path), file_extension)
            logger.info(f"文本提取完成，原始文本长度: {len(extracted_text)} 字符")
            
            # 清理文本（去除页眉页脚、水印、乱码、空白等）
            logger.info("开始清理文本内容...")
            cleaned_text = self.clean_text(extracted_text)
            logger.info(f"文本清理完成，清理后文本长度: {len(cleaned_text)} 字符")
            
            # 保存清洗后的文本到 uploads_text 目录
            logger.info("开始保存清洗后的文本文件...")
            text_file_path = self.save_cleaned_text(document_id, cleaned_text)
            logger.info(f"清洗后文本保存成功: {text_file_path}")
            
            # ⭐ 第 2 步：文本分块
            logger.info("开始对文本进行分块...")
            chunks = self.chunk_document(cleaned_text)
            logger.info(f"文本分块完成，共 {len(chunks)} 个块")
            
            # ⭐ 第 3 步：生成向量并存储
            logger.info("开始向量化和存储...")
            await self._process_and_store_chunks(document_id, chunks, file.filename)
            logger.info("向量化和存储完成")
            
            # 保存到数据库
            logger.info("开始保存文档元数据到数据库...")
            doc_metadata = repo.create(
                document_id=document_id,
                file_name=file.filename,
                file_type=file_extension,
                file_size=len(contents),
                storage_path=str(file_path),
                text_length=len(cleaned_text),
                text_file_path=text_file_path,
                metadata={
                    **(metadata or {}),
                    "text_length": len(cleaned_text),
                    "chunk_count": len(chunks)
                }
            )
            
            # 提交事务
            db_session.commit()
            logger.info(f"文档元数据保存成功: {document_id}")
            
            # 创建响应对象
            document_response = DocumentResponse(
                document_id=document_id,
                file_name=file.filename,
                status=DocumentStatus.COMPLETED,
                created_at=doc_metadata.created_at,
                updated_at=doc_metadata.updated_at,
                metadata=doc_metadata.extra_metadata  # 使用 extra_metadata
            )
            
            logger.info(f"文件处理完成: {document_id}")
            
            return DocumentUploadResult(
                document_response=document_response,
                is_duplicate=False
            )
            
        except ValueError:
            # 重新抛出业务异常
            raise
        except Exception as e:
            # 回滚事务
            db_session.rollback()
            logger.error(f"文件处理失败: {str(e)}", exc_info=True)
            
            # 删除已上传的文件
            if file_path.exists():
                try:
                    file_path.unlink()
                    logger.info(f"已删除临时文件: {file_path}")
                except Exception as cleanup_error:
                    logger.error(f"删除临时文件失败: {str(cleanup_error)}")
            
            raise Exception(f"Failed to process document: {str(e)}")
        finally:
            db_session.close()
    
    async def _process_and_store_chunks(
        self, 
        document_id: str, 
        chunks: List[str],
        file_name: str,
        use_async: bool = True  # ⭐ 新增参数：是否使用异步处理
    ):
        """
        处理并存储文本块到向量数据库（支持异步）
        
        Args:
            document_id: 文档ID
            chunks: 文本块列表
            file_name: 文件名
            use_async: 是否使用后台任务异步处理
        """
        try:
            if use_async:
                # ⭐ 使用后台任务异步处理
                from app.services.background_task_service import background_task_service
                
                logger.info(f"提交向量化后台任务: document_id={document_id}")
                
                task_id = await background_task_service.submit_task(
                    self._vectorize_and_store,
                    document_id=document_id,
                    chunks=chunks,
                    file_name=file_name
                )
                
                logger.info(f"向量化任务已提交，task_id={task_id}")
                
            else:
                # 同步处理（向后兼容）
                await self._vectorize_and_store(
                    document_id=document_id,
                    chunks=chunks,
                    file_name=file_name
                )
            
        except Exception as e:
            logger.error(f"提交向量化任务失败: {str(e)}", exc_info=True)
            # 降级为同步处理
            logger.warning("降级为同步处理向量化")
            await self._vectorize_and_store(
                document_id=document_id,
                chunks=chunks,
                file_name=file_name
            )
    
    async def _vectorize_and_store(
        self,
        document_id: str,
        chunks: List[str],
        file_name: str
    ):
        """
        执行向量化和存储（内部方法）
        
        Args:
            document_id: 文档ID
            chunks: 文本块列表
            file_name: 文件名
        """
        try:
            # 确保向量服务已初始化
            if not vector_store_service.collection:
                await vector_store_service.initialize()
            
            # 存储到向量数据库
            success = await vector_store_service.store_chunks(
                document_id=document_id,
                chunks=chunks,
                file_name=file_name
            )
            
            if not success:
                logger.warning("向量存储未成功，但继续处理")
            else:
                logger.info(f"向量化存储完成: document_id={document_id}")
            
        except Exception as e:
            logger.error(f"向量化存储失败: {str(e)}", exc_info=True)
            raise
    
    async def extract_text(self, file_path: str, file_type: str) -> str:
        """
        从不同格式的文档中提取文本
        
        Args:
            file_path: 文件路径
            file_type: 文件类型扩展名
            
        Returns:
            str: 提取的文本内容
        """
        if file_type == "pdf":
            return await self._extract_from_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            return await self._extract_from_docx(file_path)
        elif file_type in ["pptx", "ppt"]:
            return await self._extract_from_pptx(file_path)
        elif file_type in ["xlsx", "xls", "csv"]:
            return await self._extract_from_excel(file_path)
        elif file_type in ["txt", "md"]:
            return await self._extract_from_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    async def _extract_from_pdf(self, file_path: str) -> str:
        """
        从PDF提取文本
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            str: 提取的文本内容
        """
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(file_path)
            text_parts = []
            
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            return "\n".join(text_parts)
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
    
    async def _extract_from_docx(self, file_path: str) -> str:
        """
        从DOCX/DOC提取文本
        
        Args:
            file_path: Word文件路径
            
        Returns:
            str: 提取的文本内容
            
        Raises:
            Exception: 当文件损坏或格式不正确时
        """
        try:
            from docx import Document
            import zipfile
            
            # 首先验证文件是否为有效的 ZIP 格式（DOCX 本质上是 ZIP）
            logger.info(f"验证 DOCX 文件格式: {file_path}")
            try:
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    # 检查是否包含必要的 DOCX 文件
                    if '[Content_Types].xml' not in zip_ref.namelist():
                        raise ValueError("文件不是有效的 DOCX 格式：缺少 [Content_Types].xml")
                    
                    # 检查是否有损坏的文件条目
                    bad_files = [name for name in zip_ref.namelist() if 'NULL' in name.upper()]
                    if bad_files:
                        logger.warning(f"发现可疑文件条目: {bad_files}")
                        
            except zipfile.BadZipFile:
                raise ValueError("文件不是有效的 ZIP/DOCX 格式，可能已损坏")
            
            logger.info("开始解析 DOCX 文档...")
            doc = Document(file_path)
            text_parts = []
            
            # 提取段落文本
            paragraph_count = 0
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
                    paragraph_count += 1
            
            logger.info(f"提取了 {paragraph_count} 个段落")
            
            # ⭐ 提取并 OCR 处理图片（可选功能）
            try:
                image_texts = self._extract_images_from_docx(doc, file_path)
                if image_texts:
                    text_parts.append("\n\n【图片内容】\n")
                    text_parts.extend(image_texts)
                    logger.info(f"提取并OCR处理了 {len(image_texts)} 张图片")
            except Exception as e:
                logger.warning(f"图片OCR处理失败（可选功能）: {str(e)}")
            
            # 提取表格文本并转换为 Markdown 格式
            table_count = 0
            for table in doc.tables:
                table_count += 1
                markdown_table = self._convert_table_to_markdown(table)
                if markdown_table:
                    text_parts.append(markdown_table)
                    text_parts.append("")  # 表格后添加空行
            
            logger.info(f"提取了 {table_count} 个表格")
            
            if not text_parts:
                logger.warning("未提取到任何文本内容")
                return ""
            
            result = "\n".join(text_parts)
            logger.info(f"DOCX 文本提取完成，总长度: {len(result)} 字符")
            
            return result
            
        except ValueError as e:
            # 重新抛出验证错误
            logger.error(f"DOCX 文件验证失败: {str(e)}")
            raise
        except KeyError as e:
            # 处理 ZIP 条目错误
            error_msg = f"DOCX 文件结构损坏：{str(e)}"
            logger.error(error_msg)
            raise ValueError(f"文件可能已损坏或不是标准的 DOCX 格式。建议：\n"
                           f"1. 用 Microsoft Word 重新保存文件\n"
                           f"2. 尝试将文件另存为新的 .docx 文件\n"
                           f"3. 检查文件是否完整下载") from e
        except Exception as e:
            error_msg = f"DOCX extraction failed: {str(e)}"
            logger.error(error_msg, exc_info=True)
            raise Exception(error_msg)
    
    def _convert_table_to_markdown(self, table) -> str:
        """
        将 DOCX 表格转换为标准 Markdown 格式（优化版）
        
        优化点：
        1. 更好的合并单元格处理
        2. 添加表格标题（如果有）
        3. 保留单元格内的换行符为 <br>
        4. 对齐方式优化
        
        Args:
            table: python-docx Table 对象
            
        Returns:
            str: Markdown 格式的表格字符串
        """
        if not table.rows:
            return ""
        
        try:
            # 收集所有行的数据，处理合并单元格的情况
            rows_data = []
            max_cols = 0
            
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    # 清理单元格内容
                    cell_text = cell.text.strip()
                    
                    # 保留内部换行符为 HTML 换行标签（便于 LLM 理解）
                    cell_text = cell_text.replace('\n', '<br>').replace('\r', '')
                    
                    # 如果单元格为空，使用空字符串
                    cells.append(cell_text if cell_text else "")
                
                # 记录最大列数
                max_cols = max(max_cols, len(cells))
                rows_data.append(cells)
            
            if not rows_data or max_cols == 0:
                return ""
            
            # 确保所有行的列数一致，不足的补空字符串
            normalized_rows = []
            for row in rows_data:
                normalized_row = row[:max_cols]  # 截取到最大列数
                while len(normalized_row) < max_cols:
                    normalized_row.append("")
                normalized_rows.append(normalized_row)
            
            # 过滤掉全空的行
            filtered_rows = [
                row for row in normalized_rows 
                if any(cell.strip() for cell in row)
            ]
            
            if not filtered_rows:
                return ""
            
            # 计算每列的最大宽度（用于对齐）
            col_widths = [0] * max_cols
            for row in filtered_rows:
                for i, cell in enumerate(row):
                    # 计算显示宽度（排除 HTML 标签）
                    display_text = cell.replace('<br>', ' ')
                    col_widths[i] = max(col_widths[i], len(display_text))
            
            # 设置最小列宽，避免过窄
            min_col_width = 5
            col_widths = [max(w, min_col_width) for w in col_widths]
            
            # 构建 Markdown 表格
            markdown_lines = []
            
            # 表头（第一行）
            header = filtered_rows[0]
            header_line = "| " + " | ".join(
                cell.ljust(col_widths[i]) for i, cell in enumerate(header)
            ) + " |"
            markdown_lines.append(header_line)
            
            # 分隔线（居中对齐）
            separator = "|" + "|".join(
                ":" + "-" * (col_widths[i]) + ":" for i in range(max_cols)
            ) + "|"
            markdown_lines.append(separator)
            
            # 数据行（从第二行开始）
            for row in filtered_rows[1:]:
                data_line = "| " + " | ".join(
                    cell.ljust(col_widths[i]) for i, cell in enumerate(row)
                ) + " |"
                markdown_lines.append(data_line)
            
            result = "\n".join(markdown_lines)
            logger.debug(f"表格转换完成: {len(filtered_rows)} 行, {max_cols} 列")
            
            return result
            
        except Exception as e:
            logger.error(f"转换表格为 Markdown 失败: {str(e)}", exc_info=True)
            # 降级方案：返回简单的文本格式
            fallback_parts = []
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_texts.append(cell.text.strip())
                if row_texts:
                    fallback_parts.append(" | ".join(row_texts))
            return "\n".join(fallback_parts)
    
    def _convert_dataframe_to_markdown(self, df) -> str:
        """
        将 Pandas DataFrame 转换为 Markdown 表格格式
        
        Args:
            df: Pandas DataFrame 对象
            
        Returns:
            str: Markdown 格式的表格字符串
        """
        try:
            import pandas as pd
            
            if df.empty:
                return ""
            
            # 获取列名
            headers = [str(col).strip() for col in df.columns]
            num_columns = len(headers)
            
            # 将所有数据转换为字符串，并处理 NaN 值
            rows_data = []
            for _, row in df.iterrows():
                cells = []
                for cell in row:
                    if pd.notna(cell):
                        # 清理单元格内容，保留内部换行
                        cell_str = str(cell).strip().replace('\n', '<br>').replace('\r', '')
                        cells.append(cell_str)
                    else:
                        cells.append("")
                rows_data.append(cells)
            
            # 计算每列的最大宽度（用于对齐）
            col_widths = [len(h) for h in headers]
            for row in rows_data:
                for i, cell in enumerate(row):
                    if i < num_columns:
                        # 计算显示宽度（排除 HTML 标签）
                        display_text = cell.replace('<br>', ' ')
                        col_widths[i] = max(col_widths[i], len(display_text))
            
            # 设置最小列宽
            min_col_width = 5
            col_widths = [max(w, min_col_width) for w in col_widths]
            
            # 构建 Markdown 表格
            markdown_lines = []
            
            # 表头
            header_line = "| " + " | ".join(
                h.ljust(col_widths[i]) for i, h in enumerate(headers)
            ) + " |"
            markdown_lines.append(header_line)
            
            # 分隔线（居中对齐）
            separator = "|" + "|".join(
                ":" + "-" * (col_widths[i]) + ":" for i in range(num_columns)
            ) + "|"
            markdown_lines.append(separator)
            
            # 数据行
            for row in rows_data:
                data_line = "| " + " | ".join(
                    cell.ljust(col_widths[i]) if i < num_columns else ""
                    for i, cell in enumerate(row)
                ) + " |"
                markdown_lines.append(data_line)
            
            return "\n".join(markdown_lines)
            
        except Exception as e:
            logger.error(f"转换 DataFrame 为 Markdown 失败: {str(e)}", exc_info=True)
            # 降级方案
            return df.to_string(index=False)
    
    def _convert_excel_sheet_to_markdown(self, worksheet, sheet_name: str) -> str:
        """
        将 Excel 工作表转换为 Markdown 表格格式
        
        Args:
            worksheet: openpyxl Worksheet 对象
            sheet_name: 工作表名称
            
        Returns:
            str: Markdown 格式的表格字符串
        """
        try:
            # 收集所有行的数据
            rows_data = []
            max_cols = 0
            
            for row in worksheet.iter_rows(values_only=True):
                cells = [str(cell).strip() if cell is not None else "" for cell in row]
                
                # 记录最大列数
                max_cols = max(max_cols, len(cells))
                
                # 跳过全空的行
                if any(cell.strip() for cell in cells):
                    rows_data.append(cells)
            
            if not rows_data or max_cols == 0:
                return ""
            
            # 确保所有行的列数一致，不足的补空字符串
            normalized_rows = []
            for row in rows_data:
                normalized_row = row[:max_cols]  # 截取到最大列数
                while len(normalized_row) < max_cols:
                    normalized_row.append("")
                normalized_rows.append(normalized_row)
            
            # 计算每列的最大宽度（用于对齐）
            col_widths = [0] * max_cols
            for row in normalized_rows:
                for i, cell in enumerate(row):
                    col_widths[i] = max(col_widths[i], len(cell))
            
            # 构建 Markdown 表格
            markdown_lines = []
            
            # 添加工作表标题
            markdown_lines.append(f"### Sheet: {sheet_name}")
            markdown_lines.append("")
            
            # 表头（第一行）
            header = normalized_rows[0]
            header_line = "| " + " | ".join(
                cell.ljust(col_widths[i]) for i, cell in enumerate(header)
            ) + " |"
            markdown_lines.append(header_line)
            
            # 分隔线
            separator = "|" + "|".join(
                "-" * (col_widths[i] + 2) for i in range(max_cols)
            ) + "|"
            markdown_lines.append(separator)
            
            # 数据行（从第二行开始）
            for row in normalized_rows[1:]:
                data_line = "| " + " | ".join(
                    cell.ljust(col_widths[i]) for i, cell in enumerate(row)
                ) + " |"
                markdown_lines.append(data_line)
            
            result = "\n".join(markdown_lines)
            logger.debug(f"Excel工作表转换完成: {sheet_name}, {len(normalized_rows)} 行, {max_cols} 列")
            
            return result
            
        except Exception as e:
            logger.error(f"转换 Excel 工作表为 Markdown 失败: {str(e)}", exc_info=True)
            # 降级方案
            sheet_texts = []
            for row in worksheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(str(v).strip() for v in row_values):
                    sheet_texts.append(" | ".join(row_values))
            if sheet_texts:
                return f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_texts)
            return ""
    
    def _extract_images_from_docx(self, doc, file_path: str) -> List[str]:
        """
        从 DOCX 文档中提取图片并进行 OCR 识别
        
        Args:
            doc: python-docx Document 对象
            file_path: DOCX 文件路径
            
        Returns:
            List[str]: OCR 识别的文本列表
        """
        image_texts = []
        
        try:
            # 尝试导入 OCR 相关库（可选依赖）
            from PIL import Image
            import pytesseract
            import zipfile
            import io
            
            logger.info("开始提取 DOCX 中的图片...")
            
            # DOCX 本质上是 ZIP 文件，图片存储在 word/media/ 目录
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                # 查找所有图片文件
                image_files = [
                    name for name in zip_ref.namelist()
                    if name.startswith('word/media/') and 
                    name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))
                ]
                
                if not image_files:
                    logger.info("文档中没有图片")
                    return []
                
                logger.info(f"找到 {len(image_files)} 张图片")
                
                # 处理每张图片
                for i, image_name in enumerate(image_files, 1):
                    try:
                        # 读取图片数据
                        image_data = zip_ref.read(image_name)
                        
                        # 使用 Pillow 打开图片
                        image = Image.open(io.BytesIO(image_data))
                        
                        # 如果图片太大，进行缩放以提高 OCR 速度
                        max_size = 2000
                        if image.width > max_size or image.height > max_size:
                            ratio = min(max_size / image.width, max_size / image.height)
                            new_size = (int(image.width * ratio), int(image.height * ratio))
                            image = image.resize(new_size, Image.Resampling.LANCZOS)
                        
                        # 转换为灰度图（提高 OCR 准确率）
                        if image.mode != 'L':
                            image = image.convert('L')
                        
                        # 使用 Tesseract 进行 OCR
                        ocr_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                        
                        # 清理 OCR 结果
                        ocr_text = ocr_text.strip()
                        if ocr_text:
                            image_texts.append(f"[图片{i}] {ocr_text}")
                            logger.debug(f"图片 {i} OCR 完成，识别到 {len(ocr_text)} 字符")
                    
                    except Exception as e:
                        logger.warning(f"图片 {image_name} OCR 失败: {str(e)}")
                        continue
            
            return image_texts
            
        except ImportError as e:
            logger.warning(f"OCR 功能不可用（缺少依赖库）: {str(e)}")
            logger.info("如需启用 OCR 功能，请安装: pip install Pillow pytesseract")
            logger.info("同时需要安装 Tesseract OCR 引擎: https://github.com/tesseract-ocr/tesseract")
            return []
        except Exception as e:
            logger.error(f"图片提取失败: {str(e)}", exc_info=True)
            return []
    
    async def _extract_from_pptx(self, file_path: str) -> str:
        """
        从PPTX/PPT提取文本
        
        Args:
            file_path: PowerPoint文件路径
            
        Returns:
            str: 提取的文本内容
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                # 提取幻灯片标题
                if slide.shapes.title:
                    slide_texts.append(f"[Slide {slide_number}] {slide.shapes.title.text}")
                
                # 提取幻灯片内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape != slide.shapes.title:
                            slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    text_parts.append("\n".join(slide_texts))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    async def _extract_from_excel(self, file_path: str) -> str:
        """
        从Excel/CSV提取文本
        
        Args:
            file_path: Excel/CSV文件路径
            
        Returns:
            str: 提取的文本内容
        """
        try:
            import pandas as pd
            
            # 根据文件类型选择读取方式
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == ".csv":
                df = pd.read_csv(file_path)
            elif file_ext in [".xlsx", ".xls"]:
                # 读取所有sheet
                excel_file = pd.ExcelFile(file_path)
                text_parts = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    sheet_text = f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}"
                    text_parts.append(sheet_text)
                
                return "\n\n".join(text_parts)
            else:
                raise ValueError(f"Unsupported Excel format: {file_ext}")
            
            # CSV或单个sheet的处理，转换为 Markdown 表格
            return self._convert_dataframe_to_markdown(df)
            
        except ImportError:
            # 如果没有pandas，使用openpyxl
            try:
                from openpyxl import load_workbook
                
                wb = load_workbook(file_path, data_only=True)
                text_parts = []
                
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    
                    # 将 Excel 数据转换为 Markdown 表格
                    markdown_table = self._convert_excel_sheet_to_markdown(ws, sheet_name)
                    if markdown_table:
                        text_parts.append(markdown_table)
                
                return "\n\n".join(text_parts)
            except Exception as e:
                raise Exception(f"Excel extraction failed: {str(e)}")
        except Exception as e:
            raise Exception(f"Excel extraction failed: {str(e)}")
    
    async def _extract_from_text(self, file_path: str) -> str:
        """
        读取文本文件
        
        Args:
            file_path: 文本文件路径
            
        Returns:
            str: 文件内容
        """
        try:
            import chardet
            
            # 检测文件编码
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'utf-8'
            
            # 使用检测到的编码读取文件
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                return f.read()
        except Exception as e:
            # 降级方案：尝试UTF-8
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e2:
                raise Exception(f"Text file reading failed: {str(e2)}")
    
    def clean_text(self, text: str) -> str:
        """
        清理文本内容：去除页眉页脚、水印、乱码、多余空白等
        
        Args:
            text: 原始文本
            
        Returns:
            str: 清理后的文本
        """
        if not text:
            return ""
        
        # 1. 统一换行符
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        # 2. 去除常见的页眉页脚模式
        header_footer_patterns = [
            r'\n\s*Page \d+ of \d+\s*\n',  # Page 1 of 10
            r'\n\s*\d+\s*\n',  # 单独的行号
            r'\n\s*Copyright.*?\n',  # Copyright信息
            r'\n\s*Confidential.*?\n',  # 机密标记
        ]
        
        for pattern in header_footer_patterns:
            text = re.sub(pattern, '\n', text, flags=re.IGNORECASE)
        
        # 3. 去除常见的水印文本
        watermark_patterns = [
            r'DRAFT',
            r'CONFIDENTIAL',
            r'SAMPLE',
            r'WATERMARK',
        ]
        
        for pattern in watermark_patterns:
            text = re.sub(rf'\b{pattern}\b', '', text, flags=re.IGNORECASE)
        
        # 4. 去除乱码和特殊字符（保留中文、英文、数字、常见标点）
        # 保留：中英文、数字、常见标点符号、换行符、制表符
        text = re.sub(r'[^\w\s\u4e00-\u9fff\.,;:!?()\[\]{}\"\'\-—…\n\t]', ' ', text)
        
        # 5. 去除多余空白
        # 将多个空格替换为单个空格
        text = re.sub(r' +', ' ', text)
        # 将多个换行替换为最多两个换行
        text = re.sub(r'\n{3,}', '\n\n', text)
        # 去除行首行尾空白
        text = '\n'.join(line.strip() for line in text.split('\n'))
        
        # 6. 去除空行过多的部分
        lines = text.split('\n')
        filtered_lines = []
        empty_count = 0
        
        for line in lines:
            if line.strip():
                empty_count = 0
                filtered_lines.append(line)
            else:
                empty_count += 1
                if empty_count <= 2:  # 最多保留2个连续空行
                    filtered_lines.append(line)
        
        text = '\n'.join(filtered_lines)
        
        # 7. 最终清理
        text = text.strip()
        
        return text
    
    def save_cleaned_text(self, document_id: str, cleaned_text: str) -> str:
        """
        保存清洗后的文本到 uploads_text 目录
        
        Args:
            document_id: 文档UUID
            cleaned_text: 清洗后的文本内容
            
        Returns:
            str: 保存的文本文件路径
        """
        # 使用 UUID 作为文件名，扩展名为 .txt
        text_filename = f"{document_id}.txt"
        text_filepath = self.text_upload_dir / text_filename
        
        try:
            # 以 UTF-8 编码保存文本
            with open(text_filepath, 'w', encoding='utf-8') as f:
                f.write(cleaned_text)
            
            logger.info(f"清洗后的文本已保存: {text_filepath}")
            return str(text_filepath)
            
        except Exception as e:
            logger.error(f"保存清洗后文本失败: {str(e)}", exc_info=True)
            raise Exception(f"Failed to save cleaned text: {str(e)}")
    
    def chunk_document(
        self, 
        text: str, 
        chunk_size: int = None, 
        chunk_overlap: int = None
    ) -> List[str]:
        """
        将文档分割成块
        
        Args:
            text: 要分块的文本
            chunk_size: 每块的大小（字符数）
            chunk_overlap: 块之间的重叠字符数
            
        Returns:
            List[str]: 文本块列表
        """
        if chunk_size is None:
            chunk_size = settings.CHUNK_SIZE
        if chunk_overlap is None:
            chunk_overlap = settings.CHUNK_OVERLAP
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            
            # 尝试在句子边界处分割
            if end < text_length:
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                split_point = max(last_period, last_newline)
                
                if split_point > chunk_size * 0.5:  # 只有在合理位置才分割
                    chunk = chunk[:split_point + 1]
                    end = start + split_point + 1
            
            chunks.append(chunk.strip())
            start = end - chunk_overlap
        
        return chunks
    
    async def get_document(self, document_id: str) -> Optional[DocumentResponse]:
        """
        获取文档信息
        
        Args:
            document_id: 文档UUID
            
        Returns:
            Optional[DocumentResponse]: 文档响应对象
        """
        db_session = db_manager.SessionLocal()
        try:
            repo = DocumentRepository(db_session)
            doc_metadata = repo.get_by_id(document_id)
            
            if not doc_metadata:
                return None
            
            return DocumentResponse(
                document_id=doc_metadata.document_id,
                file_name=doc_metadata.file_name,
                status=DocumentStatus(doc_metadata.status),
                created_at=doc_metadata.created_at,
                updated_at=doc_metadata.updated_at,
                metadata=doc_metadata.extra_metadata  # 使用 extra_metadata
            )
        finally:
            db_session.close()
    
    async def list_documents(
        self, 
        page: int = 1, 
        page_size: int = 10,
        file_type: str = None,
        status: str = None
    ) -> List[DocumentResponse]:
        """
        列出所有文档（分页）
        
        Args:
            page: 页码
            page_size: 每页数量
            file_type: 文件类型过滤
            status: 状态过滤
            
        Returns:
            List[DocumentResponse]: 文档列表
        """
        db_session = db_manager.SessionLocal()
        try:
            repo = DocumentRepository(db_session)
            docs = repo.list_all(page=page, page_size=page_size, file_type=file_type, status=status)
            
            return [
                DocumentResponse(
                    document_id=doc.document_id,
                    file_name=doc.file_name,
                    status=DocumentStatus(doc.status),
                    created_at=doc.created_at,
                    updated_at=doc.updated_at,
                    metadata=doc.extra_metadata  # 使用 extra_metadata
                )
                for doc in docs
            ]
        finally:
            db_session.close()
    
    async def delete_document(self, document_id: str) -> bool:
        """
        删除文档及其向量数据
        
        Args:
            document_id: 文档UUID
            
        Returns:
            bool: 是否删除成功
        """
        db_session = db_manager.SessionLocal()
        try:
            repo = DocumentRepository(db_session)
            
            # 获取文档信息
            doc_metadata = repo.get_by_id(document_id)
            if not doc_metadata:
                return False
            
            # 删除文件
            file_path = Path(doc_metadata.storage_path)
            if file_path.exists():
                file_path.unlink()
            
            # 删除数据库记录
            success = repo.delete(document_id)
            
            if success:
                db_session.commit()
            
            return success
        except Exception as e:
            db_session.rollback()
            raise e
        finally:
            db_session.close()
    
    async def search_documents(self, keyword: str, page: int = 1, page_size: int = 10) -> List[DocumentResponse]:
        """
        搜索文档
        
        Args:
            keyword: 搜索关键词
            page: 页码
            page_size: 每页数量
            
        Returns:
            List[DocumentResponse]: 匹配的文档列表
        """
        db_session = db_manager.SessionLocal()
        try:
            repo = DocumentRepository(db_session)
            docs = repo.search(keyword=keyword, page=page, page_size=page_size)
            
            return [
                DocumentResponse(
                    document_id=doc.document_id,
                    file_name=doc.file_name,
                    status=DocumentStatus(doc.status),
                    created_at=doc.created_at,
                    updated_at=doc.updated_at,
                    metadata=doc.extra_metadata  # 使用 extra_metadata
                )
                for doc in docs
            ]
        finally:
            db_session.close()


# 创建全局实例
document_service = DocumentService()
